import zipfile
import os
import shutil
import sys
sys.stdout.reconfigure(encoding='utf-8', errors='replace')

pptx_path = r'C:\Users\wangh\OneDrive\Leonora\02.电风扇\健辉产品推荐.pptx'
out_dir = r'C:\BillinStone\Web\billinstone-web\public\products'
os.makedirs(out_dir, exist_ok=True)

count = 0
with zipfile.ZipFile(pptx_path, 'r') as z:
    for name in z.namelist():
        if name.startswith('ppt/media/') and any(name.lower().endswith(ext) for ext in ['.jpg', '.jpeg', '.png', '.gif', '.webp']):
            ext = os.path.splitext(name)[1].lower()
            target = os.path.join(out_dir, f'product_{count:03d}{ext}')
            with z.open(name) as src, open(target, 'wb') as dst:
                dst.write(src.read())
            count += 1

print(f'Extracted {count} images to {out_dir}')

# Also try to get text from slides
try:
    from pptx import Presentation
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides[:10]):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            print(f'\nSlide {i+1}:', ' | '.join(texts[:5]))
except ImportError:
    print('\nNote: python-pptx not installed, skipping text extraction')
